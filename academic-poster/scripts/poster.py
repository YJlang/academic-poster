"""academic-poster — a config-driven A0/A1/A2 academic poster builder.

A poster is one giant slide. This module renders a header, an optional
full-width abstract, two columns of blocks, and a footer, laying everything
out in a *flow* (sections are pushed down by measured text height) so content
never overflows into the next section. See reference/ for the why.

Usage (see examples/demo_poster.py):
    from poster import Poster
    Poster(config).save("poster.pptx")

Dependencies: python-pptx, Pillow.  Rendering a preview also needs LibreOffice.
No project-, author-, or venue-specific content lives here — pass it all via
`config`.
"""
import os
import sys

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

from pptx import Presentation
from pptx.util import Mm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData
from pptx.oxml.ns import qn

import measure

# Paper sizes in mm, portrait (w, h). Rotate for landscape.
PAPER = {"A0": (841, 1189), "A1": (594, 841), "A2": (420, 594), "A3": (297, 420)}

DEFAULT_THEME = {
    "font": "Helvetica Neue",     # font NAME written into the pptx
    "font_path": None,            # ttf path used ONLY for measuring (keep them matched!)
    "body_size": 20,
    "ink": "222A2E", "muted": "44575C",
    "primary": "2A8193", "primary_dark": "1F5C69", "primary_light": "8CC2CC",
    "accent": "D79A55",
    "box_bg": "DFEEF1", "box_bg2": "EFF6F8",
    "border": "A3CAD2", "page": "FFFFFF", "white": "FFFFFF",
    "header_text": "FFFFFF",
}
ALIGN = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
         "right": PP_ALIGN.RIGHT, "justify": PP_ALIGN.JUSTIFY}


def _rgb(hexstr):
    h = hexstr.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


class Poster:
    def __init__(self, config):
        self.cfg = config
        self.theme = {**DEFAULT_THEME, **config.get("theme", {})}
        page = config.get("page", {})
        size = page.get("size", "A1")
        w, h = PAPER.get(size, (page.get("w_mm", 594), page.get("h_mm", 841)))
        if page.get("orientation", "portrait") == "landscape":
            w, h = h, w
        self.PW, self.PH = w, h
        self.MARGIN = page.get("margin", 22)
        self.GUT = page.get("gutter", 16)
        self.font = self.theme["font"]
        self.font_path = self.theme["font_path"]
        self.BODY = self.theme["body_size"]
        self.prs = Presentation()
        self.prs.slide_width = Mm(w)
        self.prs.slide_height = Mm(h)
        self.slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

    # ---------- low-level drawing ----------
    def _set_font(self, run, size, color, bold=False):
        run.font.size = Pt(size); run.font.bold = bold
        run.font.color.rgb = _rgb(color); run.font.name = self.font
        rPr = run._r.get_or_add_rPr()
        for tag in ("a:ea", "a:cs"):
            e = rPr.find(qn(tag))
            if e is None:
                e = rPr.makeelement(qn(tag), {}); rPr.append(e)
            e.set("typeface", self.font)

    def _rect(self, x, y, w, h, color, line=None, line_w=1.0, round_=False, shadow=False, ar=0.06):
        t = MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE
        s = self.slide.shapes.add_shape(t, Mm(x), Mm(y), Mm(w), Mm(h))
        if round_:
            try: s.adjustments[0] = ar
            except Exception: pass
        s.fill.solid(); s.fill.fore_color.rgb = _rgb(color)
        if line is None: s.line.fill.background()
        else: s.line.color.rgb = _rgb(line); s.line.width = Pt(line_w)
        s.shadow.inherit = False
        if shadow: self._shadow(s)
        return s

    def _shadow(self, shape):
        spPr = shape._element.spPr
        el = spPr.makeelement(qn("a:effectLst"), {})
        sh = spPr.makeelement(qn("a:outerShdw"),
                              {"blurRad": "70000", "dist": "30000", "dir": "5400000", "rotWithShape": "0"})
        clr = spPr.makeelement(qn("a:srgbClr"), {"val": self.theme["primary_dark"]})
        clr.append(spPr.makeelement(qn("a:alpha"), {"val": "20000"}))
        sh.append(clr); el.append(sh); spPr.append(el)

    def _gradient(self, shape, c1, c2, angle=0):
        spPr = shape._element.spPr
        for tag in ("a:solidFill", "a:noFill", "a:gradFill"):
            e = spPr.find(qn(tag))
            if e is not None: spPr.remove(e)
        g = spPr.makeelement(qn("a:gradFill"), {})
        lst = spPr.makeelement(qn("a:gsLst"), {})
        for pos, col in ((0, c1), (100000, c2)):
            gs = spPr.makeelement(qn("a:gs"), {"pos": str(pos)})
            gs.append(spPr.makeelement(qn("a:srgbClr"), {"val": col.lstrip("#")}))
            lst.append(gs)
        g.append(lst)
        g.append(spPr.makeelement(qn("a:lin"), {"ang": str(int(angle * 60000)), "scaled": "1"}))
        ln = spPr.find(qn("a:ln"))
        spPr.insert(list(spPr).index(ln) if ln is not None else len(spPr), g)

    def _tb(self, x, y, w, h, anchor=MSO_ANCHOR.TOP, align="left", wrap=True):
        tb = self.slide.shapes.add_textbox(Mm(x), Mm(y), Mm(w), Mm(h))
        tf = tb.text_frame; tf.word_wrap = wrap; tf.vertical_anchor = anchor
        tf.auto_size = MSO_AUTO_SIZE.NONE      # never let the app shrink text
        for s in ("left", "right", "top", "bottom"):
            setattr(tf, "margin_" + s, Mm(0))
        tf.paragraphs[0].alignment = ALIGN[align]
        return tf

    def _para(self, tf, runs, align="left", space_after=5, line=1.27, first=False):
        p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
        p.alignment = ALIGN[align]; p.space_after = Pt(space_after); p.line_spacing = line
        for (t, sz, col, bd) in runs:
            r = p.add_run(); r.text = t; self._set_font(r, sz, col, bd)
        return p

    def _image_fit(self, path, x, y, w, h):
        from PIL import Image
        iw, ih = Image.open(path).size
        ar = iw / ih
        nw, nh = (h * ar, h) if (w / h > ar) else (w, w / ar)
        self.slide.shapes.add_picture(path, Mm(x + (w - nw) / 2), Mm(y + (h - nh) / 2), Mm(nw), Mm(nh))
        return nh

    # ---------- block renderers (return mm height consumed) ----------
    def _section_bar(self, x, y, w, title, h=16):
        bar = self._rect(x, y, w, h, self.theme["primary"], round_=True, ar=0.12)
        self._gradient(bar, self.theme["primary_dark"], self.theme["primary"], angle=0)
        tf = self._tb(x, y, w, h, anchor=MSO_ANCHOR.MIDDLE, align="center")
        self._para(tf, [("  ".join(title.upper()), 25, self.theme["white"], True)], align="center", first=True)
        return h + 6

    def _caption(self, x, y, w, num, text):
        tf = self._tb(x, y, w, 12, align="center")
        runs = [(num + " ", 15, self.theme["primary_dark"], True)] if num else []
        runs.append((text, 15, self.theme["muted"], False))
        self._para(tf, runs, align="center", first=True, line=1.1)
        return 11

    def _body(self, x, y, w, paragraphs, align="justify", gap=6.0):
        texts = [p if isinstance(p, str) else p["text"] for p in paragraphs]
        h = measure.text_height_mm(texts, w, self.BODY, self.font_path,
                                   line_spacing=1.33, para_gap_pt=gap * 0.55)
        tf = self._tb(x, y, w, h + 5)
        for i, t in enumerate(texts):
            self._para(tf, [(t, self.BODY, self.theme["ink"], False)],
                       align=align, line=1.27, space_after=gap * 0.55, first=(i == 0))
        return h

    def _bullets(self, x, y, w, items, marker="•  ", gap=3.0):
        texts = [marker + it for it in items]
        h = measure.text_height_mm(texts, w, self.BODY, self.font_path,
                                   line_spacing=1.33, para_gap_pt=gap)
        tf = self._tb(x, y, w, h + 5)
        for i, t in enumerate(texts):
            self._para(tf, [(t, self.BODY, self.theme["ink"], False)],
                       align="left", line=1.2, space_after=gap, first=(i == 0))
        return h

    def _equation(self, x, y, w, parts, number=None):
        pad = 12
        avail = (w - 2 * pad - (16 if number else 0)) * 0.90
        sz = measure.fit_one_line_size(parts, avail, self.font_path, max_size=22)
        sub = round(sz * 0.66, 1)
        fbx = sz * measure.MM_PER_PT * 1.25 + 13
        self._rect(x, y, w, fbx, self.theme["box_bg"], line=self.theme["primary"], line_w=1.2, round_=True, ar=0.07)
        ff = self._tb(x + pad, y, w - 2 * pad - (16 if number else 0), fbx, anchor=MSO_ANCHOR.MIDDLE, wrap=False)
        self._para(ff, [(t, sz if m else sub, self.theme["primary_dark"], True) for t, m in parts],
                   line=1.0, space_after=0, first=True)
        if number:
            fn = self._tb(x, y, w - pad, fbx, anchor=MSO_ANCHOR.MIDDLE, align="right")
            self._para(fn, [(number, sz, self.theme["primary_dark"], True)], align="right", first=True)
        return fbx

    def _table(self, x, y, w, block):
        data = block["data"]
        rows, cols = len(data), len(data[0])
        th = block.get("height_mm", 14 * rows)
        widths = block.get("col_widths", [1 / cols] * cols)
        gt = self.slide.shapes.add_table(rows, cols, Mm(x), Mm(y), Mm(w), Mm(th)).table
        for c in range(cols):
            gt.columns[c].width = Mm(w * widths[c])
        hl_row = block.get("highlight_row")
        for r in range(rows):
            gt.rows[r].height = Mm(th / rows)
            for c in range(cols):
                cell = gt.cell(r, c)
                cell.margin_left = Mm(3); cell.margin_right = Mm(2)
                cell.margin_top = Mm(0.3); cell.margin_bottom = Mm(0.3)
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                p = cell.text_frame.paragraphs[0]; p.alignment = ALIGN["center"]
                run = p.add_run(); run.text = str(data[r][c])
                head = (r == 0); hl = (r == hl_row)
                col = self.theme["white"] if head else (self.theme["primary_dark"] if hl else self.theme["ink"])
                self._set_font(run, 16 if head else 15, col, bold=head or hl)
                cell.fill.solid()
                if head:
                    cell.fill.fore_color.rgb = _rgb(self.theme["primary_dark"])
                elif hl:
                    cell.fill.fore_color.rgb = _rgb(self.theme["box_bg"])
                else:
                    cell.fill.fore_color.rgb = _rgb(self.theme["box_bg2"] if r % 2 else self.theme["white"])
        return th

    def _chart(self, x, y, w, block):
        cd = CategoryChartData()
        cd.categories = block["categories"]
        for s in block["series"]:
            cd.add_series(s["name"], s["values"])
        ch_h = block.get("height_mm", 80)
        gf = self.slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Mm(x), Mm(y), Mm(w), Mm(ch_h), cd)
        chart = gf.chart
        chart.has_title = False
        chart.has_legend = len(block["series"]) > 1
        if chart.has_legend:
            from pptx.enum.chart import XL_LEGEND_POSITION
            chart.legend.position = XL_LEGEND_POSITION.TOP
            chart.legend.include_in_layout = False
            chart.legend.font.size = Pt(13); chart.legend.font.name = self.font
        for si, s in enumerate(block["series"]):
            ser = chart.series[si]
            ser.format.line.fill.background()
            base = s.get("color", self.theme["primary"])
            pt_colors = s.get("point_colors")   # optional per-bar colors
            if pt_colors:
                for i, p in enumerate(ser.points):
                    p.format.fill.solid()
                    p.format.fill.fore_color.rgb = _rgb(pt_colors[i] if i < len(pt_colors) else base)
            else:
                ser.format.fill.solid(); ser.format.fill.fore_color.rgb = _rgb(base)
        plot = chart.plots[0]; plot.gap_width = block.get("gap_width", 60)
        if block.get("data_labels", True):
            plot.has_data_labels = True
            dl = plot.data_labels
            dl.number_format = block.get("label_format", "0.0"); dl.number_format_is_linked = False
            dl.position = XL_LABEL_POSITION.OUTSIDE_END
            dl.font.size = Pt(11); dl.font.bold = True; dl.font.name = self.font
            dl.font.color.rgb = _rgb(self.theme["primary_dark"])
        va = chart.value_axis
        rng = block.get("value_range", [0, 100])
        va.minimum_scale, va.maximum_scale = rng[0], rng[1]
        va.has_major_gridlines = True
        va.tick_labels.font.size = Pt(11); va.tick_labels.font.name = self.font
        va.format.line.color.rgb = _rgb(self.theme["border"])
        chart.category_axis.tick_labels.font.size = Pt(12)
        chart.category_axis.tick_labels.font.name = self.font
        return ch_h

    def _block(self, b, x, y, w):
        """Render one block, return mm consumed (incl. trailing gap)."""
        t = b["type"]
        gap = b.get("gap", 8)
        if t == "section":
            return self._section_bar(x, y, w, b["title"]) + gap - 6
        if t == "text":
            return self._body(x, y, w, b["paragraphs"], align=b.get("align", "justify")) + gap
        if t == "bullets":
            return self._bullets(x, y, w, b["items"]) + gap
        if t == "equation":
            return self._equation(x, y, w, b["parts"], b.get("number")) + gap
        if t in ("figure", "placeholder"):
            if t == "figure":
                h = self._image_fit(b["path"], x, y, w, b.get("height_mm", 120))
            else:
                ph = b.get("height_mm", 120)
                self._rect(x, y, w, ph, self.theme["box_bg"], line=self.theme["primary_light"],
                           line_w=1, round_=True, ar=0.04)
                tf = self._tb(x, y, w, ph, anchor=MSO_ANCHOR.MIDDLE, align="center")
                self._para(tf, [(b.get("label", "Figure"), 22, self.theme["primary"], True)],
                           align="center", first=True, space_after=2)
                self._para(tf, [("replace with your figure", 14, self.theme["muted"], False)], align="center")
                h = ph
            yy = y + h + 2
            if b.get("caption"):
                yy += self._caption(x, yy, w, b.get("number", ""), b["caption"])
            return (yy - y) + gap
        if t == "table":
            h = self._table(x, y, w, b)
            yy = y + h + 2
            if b.get("caption"):
                yy += self._caption(x, yy, w, b.get("number", ""), b["caption"])
            return (yy - y) + gap
        if t == "chart":
            h = self._chart(x, y, w, b)
            yy = y + h + 1
            if b.get("caption"):
                yy += self._caption(x, yy, w, b.get("number", ""), b["caption"])
            return (yy - y) + gap
        if t == "space":
            return b.get("height_mm", 8)
        raise ValueError("unknown block type: %s" % t)

    # ---------- top-level layout ----------
    def build(self):
        t, PW, PH, M, G = self.theme, self.PW, self.PH, self.MARGIN, self.GUT
        col_w = (PW - 2 * M - G) / 2
        LX, RX = M, M + col_w + G
        self._rect(0, 0, PW, PH, t["page"])
        self._rect(10, 10, PW - 20, PH - 20, t["page"], line=t["primary_light"], line_w=1.5, round_=True, ar=0.012)

        # header
        hd = self.cfg.get("header", {})
        head_h = hd.get("height_mm", 96)
        head = self._rect(M, M, PW - 2 * M, head_h, t["primary"], round_=True, ar=0.04, shadow=True)
        self._gradient(head, hd.get("color_top", t["primary"]), t["primary_dark"], angle=90)
        logo = hd.get("logo")
        if logo:
            pic = self.slide.shapes.add_picture(logo, Mm(M + 12), Mm(M), width=Mm(hd.get("logo_w", 60)))
            pic.left = Mm(M + 12); pic.top = Mm(M) + (Mm(head_h) - pic.height) // 2
        tf = self._tb(M + 10, M + 13, PW - 2 * M - 20, head_h - 24, anchor=MSO_ANCHOR.MIDDLE, align="center")
        self._para(tf, [(hd.get("title", ""), hd.get("title_size", 40), t["header_text"], True)],
                   align="center", line=1.05, space_after=10, first=True)
        if hd.get("authors"):
            self._para(tf, [(hd["authors"], 22, t["header_text"], True)], align="center", space_after=2)
        if hd.get("affiliation"):
            self._para(tf, [(hd["affiliation"], 17, t["header_text"], False)], align="center")

        y = M + head_h + 9
        # full-width abstract (optional)
        ab = self.cfg.get("abstract")
        if ab:
            y += self._section_bar(M, y, PW - 2 * M, ab.get("label", "Abstract"))
            y += self._body(M + 4, y, PW - 2 * M - 8, ab["paragraphs"]) + 12
        col_top = y

        foot_h = 24
        foot_y = PH - M - foot_h

        # two columns (flow)
        for col, x in ((self.cfg.get("left", []), LX), (self.cfg.get("right", []), RX)):
            yy = col_top
            for b in col:
                yy += self._block(b, x, yy, col_w)

        # footer
        ft = self.cfg.get("footer", {})
        foot = self._rect(M, foot_y, PW - 2 * M, foot_h, t["primary_dark"], round_=True, ar=0.18)
        self._gradient(foot, t["primary"], t["primary_dark"], angle=0)
        if ft.get("left"):
            tfl = self._tb(M + 16, foot_y, PW - 2 * M - 32, foot_h, anchor=MSO_ANCHOR.MIDDLE)
            self._para(tfl, [(ft["left"], 16, t["white"], True)], first=True)
        if ft.get("right"):
            tfr = self._tb(M + 16, foot_y, PW - 2 * M - 32, foot_h, anchor=MSO_ANCHOR.MIDDLE, align="right")
            self._para(tfr, [(ft["right"], 15, t["white"], True)], align="right", first=True)
        return self

    def save(self, path):
        self.build()
        self.prs.save(path)
        print("saved %s  (%d x %d mm)" % (path, round(self.prs.slide_width / Mm(1)),
                                          round(self.prs.slide_height / Mm(1))))
        return path
